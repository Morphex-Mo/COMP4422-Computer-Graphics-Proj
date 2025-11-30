# -*- coding: utf-8 -*-
# make_star_collector_ppt_FIXED.py   ← 改名避免覆盖旧文件
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# 配色
DARK = RGBColor(5, 15, 40)
BLUE = RGBColor(0, 120, 255)
CYAN = RGBColor(100, 220, 255)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(180, 180, 180)

def set_title(slide, text):
    title = slide.shapes.title
    title.text = text
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = BLUE
    title.text_frame.paragraphs[0].font.name = "Montserrat"
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_bullet(slide, text, size=22, color=WHITE):
    body = slide.placeholders[1].text_frame
    body.clear()
    for line in text.strip().split('\n'):
        p = body.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        if line and line[0] in "•-0-9":
            p.level = 0 if line.startswith(('•', '-')) else 1

# ====================== 正式开始 16 页 ======================

# 1 Title
s = prs.slides.add_slide(prs.slide_layouts[0])
s.background.fill.solid()
s.background.fill.fore_color.rgb = DARK
set_title(s, "Star Collector")
p = s.placeholders[1].text_frame.paragraphs[0]
p.text = "A 30-Second WebGL Storytelling Animation\n\nGroup XX • Name1, Name2, Name3"
p.font.size = Pt(28)
p.font.color.rgb = CYAN
p.alignment = PP_ALIGN.CENTER

# 2–15 内容页（严格按顺序）
content_pages = [
    ("Story & Narrative (30-second timeline)",
     "• A little fox stands on a hill at dusk, gazing at the stars\n"
     "• Three falling stars are gently collected\n"
     "• Finally, countless star particles are released into the night sky\n\n"
     "Key Moments:\n"
     "0s     → Night falls, fox breathing calmly\n"
     "11–17s → 1st star → left paw → flash & disappear\n"
     "26–32s → 2nd star → right paw → tail curls\n"
     "33s    → Camera pulls back, particles emitted upward\n\n"
     "Insert timeline screenshots here"),

    ("Scene Design",
     "• Central hill with a large tree\n"
     "• Fox positioned under the tree\n"
     "• Procedural rolling mountains (Perlin noise)\n"
     "• Sparse trees scattered on distant hills\n"
     "• Dynamic sky: sunset → deep night with visible stars\n\n"
     "Insert full-scene render here"),

    ("3D Models – Fox & Tree-Hill",
     "Fox Model\n"
     "• Source: GLTF (downloaded)\n"
     "• Triangles: 11.1k\n"
     "• Full skeletal rigging\n\n"
     "Tree-Hill Model\n"
     "• Source: GLTF\n"
     "• Triangles: 746\n"
     "• Used as central platform\n\n"
     "Insert model screenshots here"),

    ("Terrain & Stars",
     "Procedural Terrain\n"
     "• Multi-octave Perlin noise\n"
     "• Flattened center for the fox\n\n"
     "Star Objects\n"
     "• Double-layer glowing sphere\n"
     "• Core + additive glow shell\n\n"
     "Star Particles\n"
     "• Pool of 50 emissive spheres\n"
     "• Dynamic opacity & scale\n\n"
     "Insert screenshots here"),

    ("Lighting & Environment",
     "• Dynamic directional light (sun → moon)\n"
     "• Real-time day-night cycle (frozen at deep night)\n"
     "• Star self-illumination with additive glow\n"
     "• Shadow casting on fox and trees\n\n"
     "Insert day vs night comparison"),

    ("Custom Toon Shader (Fox)",
     "Features:\n"
     "• Step-wise diffuse (3 levels)\n"
     "• Hard specular highlights\n"
     "• Blue-white rim lighting\n"
     "• Shadow bands\n"
     "• Ambient term\n\n"
     "Insert 4-side comparison (standard → final toon)"),

    ("Atmospheric Scattering & Fog",
     "• Real-time Rayleigh + Mie scattering\n"
     "• Depth-based exponential fog\n"
     "• Greatly enhances depth and night mood\n\n"
     "Insert atmosphere ON/OFF comparison"),

    ("Fox Animation",
     "• Breathing → sinusoidal position offset\n"
     "• Head tilt (18–25s)\n"
     "• Tail curl after second star\n"
     "• All driven by global timer\n\n"
     "Insert animation sequence"),

    ("Star Falling & Collection",
     "• Easing (easeOutQuad)\n"
     "• Wavy drift + rotation\n"
     "• Flash → rapid fade & shrink on collection\n\n"
     "Insert falling timeline"),

    ("Particle System – Final Release",
     "• Triggered at 33s\n"
     "• 6 particles/sec, pool of 50\n"
     "• Upward acceleration + wavy drift\n"
     "• Fade-in → pulse → fade-out\n\n"
     "Insert particle fountain"),

    ("Technical Discussion – Toon Shader",
     "• Diffuse: floor(dot(N,L))\n"
     "• Specular: smoothstep threshold\n"
     "• Rim: pow(1 - dot(N,V), 2)\n"
     "• Shadow: biased shadowMap + step\n\n"
     "Insert code + diagram"),

    ("Technical Discussion – Atmosphere & Terrain",
     "• Preetham atmospheric scattering\n"
     "• 4-octave Perlin noise terrain\n"
     "• Distance-based flattening\n\n"
     "Insert heightmap"),

    ("Conclusion & Highlights",
     "Achievements\n"
     "• Complete 30s story-driven WebGL animation\n"
     "• Custom toon + atmospheric shaders\n"
     "• Sophisticated particle & animation system\n"
     "• Precise timeline synchronization\n\n"
     "Live Demo: https://your-link-here.com\n"
     "Source Code: github.com/xxx/StarCollector"),

    # 第15页：Team Contribution（之前漏掉的！）
    ("Team Contribution", 
     "Member          SID        Main Responsibilities\n"
     "───────────────────────────────────────\n"
     "Alice Wang      12345678   Fox rigging, Toon shader, Animation\n"
     "Bob Li          12345679   Terrain generation, Tree scattering\n"
     "Cathy Zhang     12345680   Star & particle system, Timeline\n"
     "All members                Atmosphere, Debug, Report & PPT\n\n"
     "(Please replace with your real info)"),
]

for title, body in content_pages:
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = DARK
    set_title(s, title)
    add_bullet(s, body, size=22 if "Team" not in title else 20, color=CYAN if "Team" in title or "Conclusion" in title else WHITE)

# 16 Q&A
s = prs.slides.add_slide(prs.slide_layouts[0])
s.background.fill.solid()
s.background.fill.fore_color.rgb = DARK
set_title(s, "Thank You!")
p = s.placeholders[1].text_frame.paragraphs[0]
p.text = "Questions?\n\nAny feedback is appreciated!"
p.font.size = Pt(36)
p.font.color.rgb = CYAN
p.alignment = PP_ALIGN.CENTER

# 保存
import os
filename = "Star_Collector_Presentation_Perfect_16pages.pptx"
prs.save(filename)

# 自动打开文件夹
print(f"成功生成完整 16 页 PPT！")
print(f"文件位置：{os.path.abspath(filename)}")
os.startfile(os.path.dirname(os.path.abspath(filename)))  # Windows
# 如果你是 Mac/Linux，把上面那行改成下面两行之一：
# subprocess.run(["open", os.path.dirname(os.path.abspath(filename))])   # Mac
# subprocess.run(["xdg-open", os.path.dirname(os.path.abspath(filename))])  # Linux